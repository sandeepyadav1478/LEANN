import argparse
import asyncio
import contextlib
import os
import time
from pathlib import Path
from typing import Any, Optional, Union

from llama_index.core import SimpleDirectoryReader
from llama_index.core.node_parser import SentenceSplitter
from tqdm import tqdm

from .api import LeannBuilder, LeannChat, LeannSearcher
from .interactive_utils import create_cli_session
from .registry import register_project_directory
from .settings import (
    resolve_anthropic_base_url,
    resolve_ollama_host,
    resolve_openai_api_key,
    resolve_openai_base_url,
)


@contextlib.contextmanager
def suppress_cpp_output(suppress: bool = True):
    """Context manager to suppress C++ stdout/stderr output from FAISS/HNSW.

    This redirects file descriptors at the OS level to suppress native C++ output
    that cannot be controlled via Python's logging framework.
    """
    if not suppress:
        yield
        return

    # Save original file descriptors
    old_stdout_fd = os.dup(1)
    old_stderr_fd = os.dup(2)

    try:
        # Open /dev/null for writing
        devnull = os.open(os.devnull, os.O_WRONLY)
        os.dup2(devnull, 1)  # Redirect stdout
        os.dup2(devnull, 2)  # Redirect stderr
        os.close(devnull)
        yield
    finally:
        # Restore original file descriptors
        os.dup2(old_stdout_fd, 1)
        os.dup2(old_stderr_fd, 2)
        os.close(old_stdout_fd)
        os.close(old_stderr_fd)


def extract_pdf_text_with_pymupdf(file_path: str) -> str | None:
    """Extract text from PDF using PyMuPDF for better quality."""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        # Check if document is encrypted and cannot be opened
        if doc.is_encrypted:
            print(f"⚠️  Skipping encrypted PDF: {file_path}")
            doc.close()
            return None
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text
    except ImportError:
        # Fallback to default reader
        return None
    except Exception as e:
        # Catch all other errors (corrupt PDFs, encoding issues, etc.)
        print(f"⚠️  Skipping problematic PDF ({type(e).__name__}): {file_path}")
        return None


def extract_pdf_text_with_pdfplumber(file_path: str) -> str | None:
    """Extract text from PDF using pdfplumber for better quality."""
    try:
        import pdfplumber

        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text
    except ImportError:
        return None
    except Exception as e:
        print(f"⚠️  Skipping problematic PDF with pdfplumber ({type(e).__name__}): {file_path}")
        return None


def extract_pdf_text_with_pypdf(file_path: str) -> str | None:
    """Extract text from PDF using pypdf as a fast fallback."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except ImportError:
        return None
    except Exception as e:
        print(f"⚠️  Skipping problematic PDF with pypdf ({type(e).__name__}): {file_path}")
        return None


def extract_pdf_docling(file_path: str) -> str | None:
    """Extract text/tables from PDF using Docling (supports OCR and complex layouts)."""
    try:
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(file_path)
        # Exporting to markdown preserves structure better for RAG
        return result.document.export_to_markdown()
    except ImportError:
        return None
    except Exception as e:
        print(f"⚠️  Skipping problematic PDF with Docling ({type(e).__name__}): {file_path}")
        return None


def extract_pdf_text(file_path: str) -> str | None:
    """Central PDF extraction with fallback chain: PyMuPDF -> pypdf -> pdfplumber -> Docling."""
    # Try PyMuPDF first (usually best layout preservation)
    text = extract_pdf_text_with_pymupdf(file_path)
    if text and text.strip():
        return text

    # Try pypdf (reliable and fast fallback)
    text = extract_pdf_text_with_pypdf(file_path)
    if text and text.strip():
        return text

    # Try pdfplumber (best for complex tables)
    text = extract_pdf_text_with_pdfplumber(file_path)
    if text and text.strip():
        return text

    # Try Docling as the ultimate fallback (OCR, complex layouts)
    text = extract_pdf_docling(file_path)
    if text and text.strip():
        return text

    return None


def extract_docx_text(file_path: str) -> str | None:
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    paragraphs.append(" | ".join(row_text))
        return "\n\n".join(paragraphs) if paragraphs else None
    except ImportError:
        return None
    except Exception as e:
        print(f"⚠️  Skipping problematic DOCX ({type(e).__name__}): {file_path}")
        return None


def extract_pptx_text(file_path: str) -> str | None:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_content.append(" | ".join(row_text))
            if slide_content:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_content))
        return "\n\n".join(slides_text) if slides_text else None
    except ImportError:
        return None
    except Exception as e:
        print(f"⚠️  Skipping problematic PPTX ({type(e).__name__}): {file_path}")
        return None


def extract_xlsx_text(file_path: str) -> str | None:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    rows.append(" | ".join(row_values))
            if rows:
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets_text) if sheets_text else None
    except ImportError:
        return None
    except Exception as e:
        print(f"⚠️  Skipping problematic XLSX ({type(e).__name__}): {file_path}")
        return None


def extract_mm_text(file_path: str) -> str | None:
    """Extract text from FreeMind/Freeplane .mm files (XML-based mindmaps)."""
    try:
        import xml.etree.ElementTree as ET

        tree = ET.parse(file_path)
        root = tree.getroot()
        extracted_texts = []

        def traverse_node(node_element, level=0):
            indent = "  " * level
            text_attr = node_element.get("TEXT")
            if text_attr:
                extracted_texts.append(f"{indent}- {text_attr}")

            # Extract text from <note> child element (FreeMind notes)
            note_element = node_element.find("note")
            if note_element is not None and note_element.text:
                note_text = note_element.text.strip()
                if note_text:
                    extracted_texts.append(f"{indent}  (Note: {note_text})")

            # For Freeplane, notes can be in <richcontent TYPE="NOTE">
            for rich in node_element.findall("richcontent"):
                if rich.get("TYPE") == "NOTE":
                    # Simple text extraction from rich content
                    rich_text = "".join(rich.itertext()).strip()
                    if rich_text:
                        extracted_texts.append(f"{indent}  (Note: {rich_text})")

            # Recursively traverse child nodes
            for child_node in node_element.findall("node"):
                traverse_node(child_node, level + 1)

        # Start traversal from any top-level node elements
        for node_child in root.findall("node"):
            traverse_node(node_child, level=0)

        return "\n".join(extracted_texts) if extracted_texts else None
    except Exception as e:
        print(f"⚠️  Skipping problematic MM ({type(e).__name__}): {file_path} - {e}")
        return None


class LeannCLI:
    def __init__(self):
        # Central storage logic
        leann_home_env = os.environ.get("LEANN_HOME")
        if leann_home_env:
            # If LEANN_HOME is set, use it as the base (global/central storage)
            self.leann_home = Path(leann_home_env).expanduser().resolve()
        else:
            # Otherwise, use project-local .leann directory (like .git)
            self.leann_home = Path.cwd() / ".leann"

        self.indexes_dir = self.leann_home / "indexes"
        self.indexes_dir.mkdir(parents=True, exist_ok=True)

        # Default parser for documents
        self.node_parser = SentenceSplitter(
            chunk_size=256, chunk_overlap=128, separator=" ", paragraph_separator="\n\n"
        )

        # Code-optimized parser
        self.code_parser = SentenceSplitter(
            chunk_size=512,  # Larger chunks for code context
            chunk_overlap=50,  # Less overlap to preserve function boundaries
            separator="\n",  # Split by lines for code
            paragraph_separator="\n\n",  # Preserve logical code blocks
        )

    def get_index_path(self, index_name: str) -> str:
        index_dir = self.indexes_dir / index_name
        return str(index_dir / "documents.leann")

    def index_exists(self, index_name: str) -> bool:
        index_dir = self.indexes_dir / index_name
        meta_file = index_dir / "documents.leann.meta.json"
        return meta_file.exists()

    def create_parser(self) -> argparse.ArgumentParser:
        parser = argparse.ArgumentParser(
            prog="leann",
            description="The smallest vector index in the world. RAG Everything with LEANN!",
            formatter_class=argparse.RawDescriptionHelpFormatter,
            epilog="""
Examples:
  leann build my-docs --docs ./documents                                  # Build index from directory
  leann build my-code --docs ./src ./tests ./config                      # Build index from multiple directories
  leann build my-files --docs ./file1.py ./file2.txt ./docs/             # Build index from files and directories
  leann build my-mixed --docs ./readme.md ./src/ ./config.json           # Build index from mixed files/dirs
  leann build my-ppts --docs ./ --file-types .pptx,.pdf                  # Index only PowerPoint and PDF files
  leann search my-docs "query"                                           # Search in my-docs index
  leann ask my-docs "question"                                           # Ask my-docs index
  leann react my-docs "complex question"                                 # Use ReAct agent for multiturn retrieval
  leann list                                                             # List all stored indexes
  leann remove my-docs                                                   # Remove an index (local first, then global)
            """,
        )

        # Global verbosity options
        verbosity_group = parser.add_mutually_exclusive_group()
        verbosity_group.add_argument(
            "-v",
            "--verbose",
            action="store_true",
            help="Show detailed output including C++ backend logs from FAISS/HNSW",
        )
        verbosity_group.add_argument(
            "-q",
            "--quiet",
            action="store_true",
            help="Suppress all non-essential output (default behavior)",
        )

        subparsers = parser.add_subparsers(dest="command", help="Available commands")

        # Build command
        build_parser = subparsers.add_parser("build", help="Build document index")
        build_parser.add_argument(
            "index_name", nargs="?", help="Index name (default: current directory name)"
        )
        build_parser.add_argument(
            "--docs",
            type=str,
            nargs="+",
            default=["."],
            help="Documents directories and/or files (default: current directory)",
        )
        build_parser.add_argument(
            "--backend-name",
            type=str,
            default="hnsw",
            choices=["hnsw", "diskann"],
            help="Backend to use (default: hnsw)",
        )
        build_parser.add_argument(
            "--embedding-model",
            type=str,
            default="facebook/contriever",
            help="Embedding model (default: facebook/contriever)",
        )
        build_parser.add_argument(
            "--embedding-mode",
            type=str,
            default="sentence-transformers",
            choices=["sentence-transformers", "openai", "mlx", "ollama"],
            help="Embedding backend mode (default: sentence-transformers)",
        )
        build_parser.add_argument(
            "--embedding-host",
            type=str,
            default=None,
            help="Override Ollama-compatible embedding host",
        )
        build_parser.add_argument(
            "--embedding-api-base",
            type=str,
            default=None,
            help="Base URL for OpenAI-compatible embedding services",
        )
        build_parser.add_argument(
            "--embedding-api-key",
            type=str,
            default=None,
            help="API key for embedding service (defaults to OPENAI_API_KEY)",
        )
        build_parser.add_argument(
            "--embedding-prompt-template",
            type=str,
            default=None,
            help="Prompt template to prepend to all texts for embedding (e.g., 'query: ' for search)",
        )
        build_parser.add_argument(
            "--query-prompt-template",
            type=str,
            default=None,
            help="Prompt template for queries (different from build template for task-specific models)",
        )
        build_parser.add_argument(
            "--force", "-f", action="store_true", help="Force rebuild existing index"
        )
        build_parser.add_argument(
            "--graph-degree", type=int, default=32, help="Graph degree (default: 32)"
        )
        build_parser.add_argument(
            "--complexity", type=int, default=64, help="Build complexity (default: 64)"
        )
        build_parser.add_argument("--num-threads", type=int, default=1)
        build_parser.add_argument(
            "--compact",
            action=argparse.BooleanOptionalAction,
            default=True,
            help="Use compact storage (default: true). Must be `no-compact` for `no-recompute` build.",
        )
        build_parser.add_argument(
            "--recompute",
            action=argparse.BooleanOptionalAction,
            default=True,
            help="Enable recomputation (default: true)",
        )
        build_parser.add_argument(
            "--file-types",
            type=str,
            help="Comma-separated list of file extensions to include (e.g., '.txt,.pdf,.pptx'). If not specified, uses default supported types.",
        )
        build_parser.add_argument(
            "--include-hidden",
            action=argparse.BooleanOptionalAction,
            default=False,
            help="Include hidden files and directories (paths starting with '.') during indexing (default: false)",
        )
        build_parser.add_argument(
            "--doc-chunk-size",
            type=int,
            default=256,
            help="Document chunk size in TOKENS (default: 256). Final chunks may be larger due to overlap. For 512 token models: recommended 350 tokens (350 + 128 overlap = 478 max)",
        )
        build_parser.add_argument(
            "--doc-chunk-overlap",
            type=int,
            default=128,
            help="Document chunk overlap in TOKENS (default: 128). Added to chunk size, not included in it",
        )
        build_parser.add_argument(
            "--code-chunk-size",
            type=int,
            default=512,
            help="Code chunk size in TOKENS (default: 512). Final chunks may be larger due to overlap. For 512 token models: recommended 400 tokens (400 + 50 overlap = 450 max)",
        )
        build_parser.add_argument(
            "--code-chunk-overlap",
            type=int,
            default=50,
            help="Code chunk overlap in TOKENS (default: 50). Added to chunk size, not included in it",
        )
        build_parser.add_argument(
            "--use-ast-chunking",
            action="store_true",
            help="Enable AST-aware chunking for code files (requires astchunk)",
        )
        build_parser.add_argument(
            "--ast-chunk-size",
            type=int,
            default=300,
            help="AST chunk size in CHARACTERS (non-whitespace) (default: 300). Final chunks may be larger due to overlap and expansion. For 512 token models: recommended 300 chars (300 + 64 overlap ~= 480 tokens)",
        )
        build_parser.add_argument(
            "--ast-chunk-overlap",
            type=int,
            default=64,
            help="AST chunk overlap in CHARACTERS (default: 64). Added to chunk size, not included in it. ~1.2 tokens per character for code",
        )
        build_parser.add_argument(
            "--ast-fallback-traditional",
            action="store_true",
            default=True,
            help="Fall back to traditional chunking if AST chunking fails (default: True)",
        )

        # Search command
        search_parser = subparsers.add_parser("search", help="Search documents")
        search_parser.add_argument("index_name", help="Index name")
        search_parser.add_argument("query", help="Search query")
        search_parser.add_argument(
            "--top-k", type=int, default=5, help="Number of results (default: 5)"
        )
        search_parser.add_argument(
            "--complexity", type=int, default=64, help="Search complexity (default: 64)"
        )
        search_parser.add_argument("--beam-width", type=int, default=1)
        search_parser.add_argument("--prune-ratio", type=float, default=0.0)
        search_parser.add_argument(
            "--recompute",
            dest="recompute_embeddings",
            action=argparse.BooleanOptionalAction,
            default=True,
            help="Enable/disable embedding recomputation (default: enabled). Should not do a `no-recompute` search in a `recompute` build.",
        )
        search_parser.add_argument(
            "--pruning-strategy",
            choices=["global", "local", "proportional"],
            default="global",
            help="Pruning strategy (default: global)",
        )
        search_parser.add_argument(
            "--non-interactive",
            action="store_true",
            help="Non-interactive mode: automatically select index without prompting",
        )
        search_parser.add_argument(
            "--show-metadata",
            action="store_true",
            help="Display file paths and metadata in search results",
        )
        search_parser.add_argument(
            "--embedding-prompt-template",
            type=str,
            default=None,
            help="Prompt template to prepend to query for embedding (e.g., 'query: ' for search)",
        )

        # Ask command
        ask_parser = subparsers.add_parser("ask", help="Ask questions")
        ask_parser.add_argument("index_name", help="Index name")
        ask_parser.add_argument(
            "query",
            nargs="?",
            help="Question to ask (omit for prompt or when using --interactive)",
        )
        ask_parser.add_argument(
            "--llm",
            type=str,
            default="ollama",
            choices=["simulated", "ollama", "hf", "openai", "anthropic"],
            help="LLM provider (default: ollama)",
        )
        ask_parser.add_argument(
            "--model", type=str, default="qwen3:8b", help="Model name (default: qwen3:8b)"
        )
        ask_parser.add_argument(
            "--host",
            type=str,
            default=None,
            help="Override Ollama-compatible host (defaults to LEANN_OLLAMA_HOST/OLLAMA_HOST)",
        )
        ask_parser.add_argument(
            "--interactive", "-i", action="store_true", help="Interactive chat mode"
        )
        ask_parser.add_argument(
            "--top-k", type=int, default=20, help="Retrieval count (default: 20)"
        )
        ask_parser.add_argument("--complexity", type=int, default=32)
        ask_parser.add_argument("--beam-width", type=int, default=1)
        ask_parser.add_argument("--prune-ratio", type=float, default=0.0)
        ask_parser.add_argument(
            "--recompute",
            dest="recompute_embeddings",
            action=argparse.BooleanOptionalAction,
            default=True,
            help="Enable/disable embedding recomputation during ask (default: enabled)",
        )
        ask_parser.add_argument(
            "--pruning-strategy",
            choices=["global", "local", "proportional"],
            default="global",
        )
        ask_parser.add_argument(
            "--thinking-budget",
            type=str,
            choices=["low", "medium", "high"],
            default=None,
            help="Thinking budget for reasoning models (low/medium/high). Supported by GPT-Oss:20b and other reasoning models.",
        )
        ask_parser.add_argument(
            "--api-base",
            type=str,
            default=None,
            help="Base URL for OpenAI-compatible APIs (e.g., http://localhost:10000/v1)",
        )
        ask_parser.add_argument(
            "--api-key",
            type=str,
            default=None,
            help="API key for cloud LLM providers (OpenAI, Anthropic)",
        )

        # React command (multiturn retrieval agent)
        react_parser = subparsers.add_parser(
            "react", help="Use ReAct agent for multiturn retrieval and reasoning"
        )
        react_parser.add_argument("index_name", help="Index name")
        react_parser.add_argument("query", help="Question to research")
        react_parser.add_argument(
            "--llm",
            type=str,
            default="ollama",
            choices=["simulated", "ollama", "hf", "openai", "anthropic"],
            help="LLM provider (default: ollama)",
        )
        react_parser.add_argument(
            "--model", type=str, default="qwen3:8b", help="Model name (default: qwen3:8b)"
        )
        react_parser.add_argument(
            "--host",
            type=str,
            default=None,
            help="Override Ollama-compatible host (defaults to LEANN_OLLAMA_HOST/OLLAMA_HOST)",
        )
        react_parser.add_argument(
            "--top-k", type=int, default=5, help="Number of results per search (default: 5)"
        )
        react_parser.add_argument(
            "--max-iterations",
            type=int,
            default=5,
            help="Maximum number of search iterations (default: 5)",
        )
        react_parser.add_argument(
            "--api-base",
            type=str,
            default=None,
            help="Base URL for OpenAI-compatible APIs (e.g., http://localhost:10000/v1)",
        )
        react_parser.add_argument(
            "--api-key",
            type=str,
            default=None,
            help="API key for cloud LLM providers (OpenAI, Anthropic)",
        )

        # List command
        subparsers.add_parser("list", help="List all indexes")

        # Remove command
        remove_parser = subparsers.add_parser("remove", help="Remove an index")
        remove_parser.add_argument("index_name", help="Index name to remove")
        remove_parser.add_argument(
            "--force", "-f", action="store_true", help="Force removal without confirmation"
        )

        # Serve command (HTTP API server)
        serve_parser = subparsers.add_parser(
            "serve", help="Start HTTP API server for LEANN vector DB"
        )
        serve_parser.add_argument(
            "--host", type=str, default=None, help="Host to bind to (default: 0.0.0.0)"
        )
        serve_parser.add_argument(
            "--port", type=int, default=None, help="Port to bind to (default: 8000)"
        )

        # Browser Index Command
        browser_parser = subparsers.add_parser("index-browser", help="Index browser history")
        browser_parser.add_argument(
            "browser_type", choices=["chrome", "brave"], help="Type of browser"
        )
        browser_parser.add_argument(
            "--profile", type=str, default="Default", help="Profile name (default: Default)"
        )
        browser_parser.add_argument(
            "--index-name", type=str, default=None, help="Custom index name"
        )
        browser_parser.add_argument(
            "--max-items", type=int, default=1000, help="Max history items to index"
        )

        # Email indexing command
        email_parser = subparsers.add_parser("index-email", help="Index Apple Mail emails")
        email_parser.add_argument(
            "index_name", nargs="?", default="apple-mail", help="Index name (default: apple-mail)"
        )
        email_parser.add_argument(
            "--max-items", type=int, default=2000, help="Max emails to index (default: 2000)"
        )

        # Calendar indexing command
        calendar_parser = subparsers.add_parser(
            "index-calendar", help="Index Apple Calendar events"
        )
        calendar_parser.add_argument(
            "index_name",
            nargs="?",
            default="apple-calendar",
            help="Index name (default: apple-calendar)",
        )
        calendar_parser.add_argument(
            "--max-items", type=int, default=1000, help="Max events to index (default: 1000)"
        )

        # WeChat indexing command
        wechat_parser = subparsers.add_parser("index-wechat", help="Index WeChat chat history")
        wechat_parser.add_argument(
            "index_name", nargs="?", default="wechat", help="Index name (default: wechat)"
        )
        wechat_parser.add_argument(
            "--export-dir",
            type=str,
            default="./wechat_export",
            help="Directory for WeChat exports",
        )
        wechat_parser.add_argument(
            "--max-items", type=int, default=1000, help="Max chat entries to index"
        )

        # iMessage indexing command
        imessage_parser = subparsers.add_parser("index-imessage", help="Index iMessage history")
        imessage_parser.add_argument(
            "index_name", nargs="?", default="imessage", help="Index name (default: imessage)"
        )
        imessage_parser.add_argument("--db-path", type=str, help="Custom chat.db path")
        imessage_parser.add_argument(
            "--max-items", type=int, default=2000, help="Max messages to index"
        )

        # Slack indexing command
        slack_parser = subparsers.add_parser("index-slack", help="Index Slack workspace via MCP")
        slack_parser.add_argument(
            "index_name", nargs="?", default="slack", help="Index name (default: slack)"
        )
        slack_parser.add_argument(
            "--mcp-server",
            type=str,
            required=True,
            help="MCP server command (e.g. 'slack-mcp-server')",
        )
        slack_parser.add_argument("--workspace", type=str, help="Slack workspace name")
        slack_parser.add_argument("--channels", nargs="+", help="Specific channels to index")

        # ChatGPT indexing command
        chatgpt_parser = subparsers.add_parser("index-chatgpt", help="Index ChatGPT export")
        chatgpt_parser.add_argument(
            "index_name", nargs="?", default="chatgpt", help="Index name (default: chatgpt)"
        )
        chatgpt_parser.add_argument(
            "--export-path", type=str, required=True, help="Path to export file"
        )

        # Claude indexing command
        claude_parser = subparsers.add_parser("index-claude", help="Index Claude export")
        claude_parser.add_argument(
            "index_name", nargs="?", default="claude", help="Index name (default: claude)"
        )
        claude_parser.add_argument(
            "--export-path", type=str, required=True, help="Path to export file"
        )

        return parser

    def register_project_dir(self):
        """Register current project directory in global registry"""
        register_project_directory()

    def _build_gitignore_parser(self, docs_dir: str):
        """Build gitignore parser using gitignore-parser library."""
        from gitignore_parser import parse_gitignore

        # Try to parse the root .gitignore
        gitignore_path = Path(docs_dir) / ".gitignore"

        if gitignore_path.exists():
            try:
                # gitignore-parser automatically handles all subdirectory .gitignore files!
                matches = parse_gitignore(str(gitignore_path))
                print(f"📋 Loaded .gitignore from {docs_dir} (includes all subdirectories)")
                return matches
            except Exception as e:
                print(f"Warning: Could not parse .gitignore: {e}")
        else:
            print("📋 No .gitignore found")

        # Fallback: basic pattern matching for essential files
        essential_patterns = {".git", ".DS_Store", "__pycache__", "node_modules", ".venv", "venv"}

        def basic_matches(file_path):
            path_parts = Path(file_path).parts
            return any(part in essential_patterns for part in path_parts)

        return basic_matches

    def _should_exclude_file(self, file_path: Path, gitignore_matches) -> bool:
        """Check if a file should be excluded using gitignore parser.

        Always match against absolute, posix-style paths for consistency with
        gitignore_parser expectations.
        """
        try:
            absolute_path = file_path.resolve()
        except Exception:
            absolute_path = Path(str(file_path))
        return gitignore_matches(absolute_path.as_posix())

    def _is_git_submodule(self, path: Path) -> bool:
        """Check if a path is a git submodule."""
        try:
            # Find the git repo root
            current_dir = Path.cwd()
            while current_dir != current_dir.parent:
                if (current_dir / ".git").exists():
                    gitmodules_path = current_dir / ".gitmodules"
                    if gitmodules_path.exists():
                        # Read .gitmodules to check if this path is a submodule
                        gitmodules_content = gitmodules_path.read_text()
                        # Convert path to relative to git root
                        try:
                            relative_path = path.resolve().relative_to(current_dir)
                            # Check if this path appears in .gitmodules
                            return f"path = {relative_path}" in gitmodules_content
                        except ValueError:
                            # Path is not under git root
                            return False
                    break
                current_dir = current_dir.parent
            return False
        except Exception:
            # If anything goes wrong, assume it's not a submodule
            return False

    def list_indexes(self):
        # Get all project directories with .leann
        global_registry = Path.home() / ".leann" / "projects.json"
        all_projects = []

        if global_registry.exists():
            try:
                import json

                with open(global_registry) as f:
                    all_projects = json.load(f)
            except Exception:
                pass

        # Filter to only existing directories with .leann
        valid_projects = []
        for project_dir in all_projects:
            project_path = Path(project_dir)
            if project_path.exists() and (project_path / ".leann" / "indexes").exists():
                valid_projects.append(project_path)

        # Add current project if it has .leann but not in registry
        current_path = Path.cwd()
        if (current_path / ".leann" / "indexes").exists() and current_path not in valid_projects:
            valid_projects.append(current_path)

        # Separate current and other projects
        other_projects = []

        for project_path in valid_projects:
            if project_path != current_path:
                other_projects.append(project_path)

        print("📚 LEANN Indexes")
        print("=" * 50)

        total_indexes = 0
        current_indexes_count = 0
        size_cache = {}

        # Show current storage header
        leann_home_env = os.environ.get("LEANN_HOME")
        if leann_home_env:
            print("\n🏠 Central Index Storage (LEANN_HOME)")
            print(f"   {self.leann_home}")
        else:
            print("\n🏠 Current Project (Local)")
            print(f"   {current_path}")
        print("   " + "─" * 45)

        # Optimization: Only include central indexes in the current project listing
        current_indexes = self._discover_indexes_in_project(
            current_path, exclude_dirs=other_projects, include_central=True, size_cache=size_cache
        )
        if current_indexes:
            for idx in current_indexes:
                total_indexes += 1
                current_indexes_count += 1
                type_icon = "📁" if idx["type"] == "cli" else "📄"
                print(f"   {current_indexes_count}. {type_icon} {idx['name']} {idx['status']}")
                if idx["size_mb"] > 0:
                    print(f"      📦 Size: {idx['size_mb']:.1f} MB")
        else:
            print("   📭 No indexes in current project")

        # Show other projects (reference information)
        if other_projects:
            print("\n\n🗂️  Other Projects")
            print("   " + "─" * 45)

            for project_path in other_projects:
                # Optimization: Do NOT include central indexes in other project listings
                # This prevents massive redundancy and saves time
                project_indexes = self._discover_indexes_in_project(
                    project_path, include_central=False, size_cache=size_cache
                )
                if not project_indexes:
                    continue

                print(f"\n   📂 {project_path.name}")
                print(f"      {project_path}")

                for idx in project_indexes:
                    total_indexes += 1
                    type_icon = "📁" if idx["type"] == "cli" else "📄"
                    print(f"      • {type_icon} {idx['name']} {idx['status']}")
                    if idx["size_mb"] > 0:
                        print(f"        📦 {idx['size_mb']:.1f} MB")

        # Summary and usage info
        print("\n" + "=" * 50)
        if total_indexes == 0:
            print("💡 Get started:")
            print("   leann build my-docs --docs ./documents")
        else:
            print(f"📊 Total: {total_indexes} indexes")

            if current_indexes_count > 0:
                print("\n💫 Quick start:")
                # Get first index from current project for example
                example_idx = current_indexes[0]
                example_name = example_idx["name"]
                print(f'   leann search {example_name} "your query"')
                print(f"   leann ask {example_name} --interactive")
            else:
                print("\n💡 Create your first index:")
                print("   leann build my-docs --docs ./documents")

    def _discover_indexes_in_project(
        self,
        project_path: Path,
        exclude_dirs: Optional[list[Path]] = None,
        include_central: bool = True,
        size_cache: Optional[dict[str, float]] = None,
    ):
        """Discover all indexes in a project directory (both CLI and apps formats)

        include_central: If True, includes the central indexes directory (~/.leann/indexes)
        exclude_dirs: when provided, skip any APP-format index files that are
        located under these directories. This prevents duplicates when the
        current project is a parent directory of other registered projects.
        size_cache: Cache for index sizes to avoid redundant stat calls
        """
        indexes = []
        exclude_dirs = exclude_dirs or []
        size_cache = size_cache if size_cache is not None else {}

        # normalize to resolved paths once for comparison
        try:
            exclude_dirs_resolved = [p.resolve() for p in exclude_dirs]
        except Exception:
            exclude_dirs_resolved = exclude_dirs

        # 1. CLI format: .leann/indexes/index_name/ (Check both central and local)
        dirs_to_check = []
        if include_central:
            dirs_to_check.append(self.indexes_dir)

        # If project_path has a local .leann/indexes, also check it
        try:
            local_indexes_dir = project_path / ".leann" / "indexes"
            if local_indexes_dir.exists():
                # Avoid duplicate scan if local matches central or is inside it
                local_resolved = local_indexes_dir.resolve()
                central_resolved = self.indexes_dir.resolve()
                try:
                    if local_resolved != central_resolved and not local_resolved.is_relative_to(
                        central_resolved
                    ):
                        dirs_to_check.append(local_indexes_dir)
                except AttributeError:  # Python < 3.9 fallback
                    if str(local_resolved) != str(central_resolved) and str(
                        central_resolved
                    ) not in str(local_resolved):
                        dirs_to_check.append(local_indexes_dir)
        except (OSError, PermissionError):
            pass

        for cli_indexes_dir in dirs_to_check:
            try:
                if cli_indexes_dir.exists():
                    for index_dir in cli_indexes_dir.iterdir():
                        if index_dir.is_dir():
                            # Cache key for this index
                            cache_key = str(index_dir.resolve())

                            meta_file = index_dir / "documents.leann.meta.json"
                            status = "✅" if meta_file.exists() else "❌"

                            size_mb = 0
                            if meta_file.exists():
                                if cache_key in size_cache:
                                    size_mb = size_cache[cache_key]
                                else:
                                    try:
                                        # Fast size calculation (top-level files only)
                                        size_mb = sum(
                                            f.stat().st_size
                                            for f in index_dir.iterdir()
                                            if f.is_file()
                                        ) / (1024 * 1024)
                                        size_cache[cache_key] = size_mb
                                    except (OSError, PermissionError):
                                        pass

                            # Avoid duplicates if checking multiple directories
                            if any(idx["name"] == index_dir.name for idx in indexes):
                                continue

                            indexes.append(
                                {
                                    "name": index_dir.name,
                                    "type": "cli",
                                    "status": status,
                                    "size_mb": size_mb,
                                    "path": index_dir,
                                }
                            )
            except (OSError, PermissionError, FileNotFoundError):
                continue

        # 2. Apps format: *.leann.meta.json files in the project root ONLY
        # We avoid recursive scanning to ensure high performance
        meta_files = []
        try:
            # Level 0 ONLY (Project Root)
            meta_files.extend(project_path.glob("*.leann.meta.json"))
        except (OSError, PermissionError):
            pass

        exclusion_indexes_dir = self.indexes_dir
        for meta_file in meta_files:
            try:
                if meta_file.is_file():
                    # Skip CLI-built indexes
                    try:
                        if (
                            exclusion_indexes_dir.exists()
                            and exclusion_indexes_dir in meta_file.parents
                        ):
                            continue
                    except Exception:
                        pass

                    # Skip meta files that live under excluded directories
                    try:
                        meta_parent_resolved = meta_file.parent.resolve()
                        if any(
                            meta_parent_resolved.is_relative_to(ex_dir)
                            for ex_dir in exclude_dirs_resolved
                        ):
                            continue
                    except Exception:
                        pass

                    display_name = meta_file.parent.name
                    file_base = meta_file.name.replace(".leann.meta.json", "")
                    status = "✅"

                    # Cache key for app index
                    cache_key = f"app:{meta_file.resolve()}"

                    size_mb = 0
                    if cache_key in size_cache:
                        size_mb = size_cache[cache_key]
                    else:
                        try:
                            index_dir = meta_file.parent
                            # Optimization: Avoid glob in large directories
                            related_extensions = [
                                "",
                                ".passages.jsonl",
                                ".passages.idx",
                                ".mapping.json",
                            ]
                            for rel_ext in related_extensions:
                                rel_path = index_dir / f"{file_base}.leann{rel_ext}"
                                try:
                                    if rel_path.is_file():
                                        size_mb += rel_path.stat().st_size / (1024 * 1024)
                                except (OSError, PermissionError):
                                    continue
                            size_cache[cache_key] = size_mb
                        except (OSError, PermissionError):
                            pass

                    # Avoid duplicates
                    if any(idx["name"] == display_name and idx["type"] == "app" for idx in indexes):
                        continue

                    indexes.append(
                        {
                            "name": display_name,
                            "type": "app",
                            "status": status,
                            "size_mb": size_mb,
                            "path": meta_file,
                        }
                    )
            except (OSError, PermissionError, FileNotFoundError):
                continue

        return indexes

    def remove_index(self, index_name: str, force: bool = False):
        """Safely remove an index - always show all matches for transparency"""

        # Always do a comprehensive search for safety
        print(f"🔍 Searching for all indexes named '{index_name}'...")
        all_matches = self._find_all_matching_indexes(index_name)

        if not all_matches:
            print(f"❌ Index '{index_name}' not found in any project.")
            return False

        if len(all_matches) == 1:
            return self._remove_single_match(all_matches[0], index_name, force)
        else:
            return self._remove_from_multiple_matches(all_matches, index_name, force)

    def _find_all_matching_indexes(self, index_name: str):
        """Find all indexes with the given name across all projects"""
        matches = []

        # Optimization: use the same discovery logic as list_indexes but filter by name
        # First, find all projects
        global_registry = Path.home() / ".leann" / "projects.json"
        all_projects = []
        if global_registry.exists():
            try:
                import json

                with open(global_registry) as f:
                    all_projects = json.load(f)
            except Exception:
                pass

        current_path = Path.cwd()
        project_paths = [Path(p) for p in all_projects if Path(p).exists()]
        if current_path not in project_paths:
            project_paths.append(current_path)

        # Track seen paths to avoid duplicates
        seen_dirs = set()

        for project_path in project_paths:
            # Discover indexes in this project using the optimized method
            # We exclude central from others to avoid massive redundancy
            is_current = project_path == current_path
            discovered = self._discover_indexes_in_project(project_path, include_central=is_current)

            for idx in discovered:
                if idx["name"] == index_name:
                    # Resolve path for deduplication
                    path_resolved = idx["path"].resolve()
                    if path_resolved in seen_dirs:
                        continue
                    seen_dirs.add(path_resolved)

                    matches.append(
                        {
                            "project_path": project_path,
                            "index_dir": idx["path"]
                            if idx["type"] == "cli"
                            else idx["path"].parent,
                            "meta_file": idx["path"]
                            if idx["type"] == "app"
                            else (idx["path"] / "documents.leann.meta.json"),
                            "is_current": is_current,
                            "kind": idx["type"],
                            "display_name": idx["name"],
                            "file_base": idx["path"].name if idx["type"] == "app" else "documents",
                            "path": idx["path"],
                        }
                    )

        # Sort: current project first, then by project name
        matches.sort(key=lambda x: (not x["is_current"], x["project_path"].name))
        return matches

    def _remove_single_match(self, match, index_name: str, force: bool):
        """Handle removal when only one match is found"""
        project_path = match["project_path"]
        is_current = match["is_current"]
        kind = match.get("kind", "cli")

        if is_current:
            location_info = "current project"
            emoji = "🏠"
        else:
            location_info = f"other project '{project_path.name}'"
            emoji = "📂"

        print(f"✅ Found 1 index named '{index_name}':")
        print(f"   {emoji} Location: {location_info}")
        if kind == "cli":
            print(f"   📍 Path: {match['path']}")
        else:
            print(f"   📍 Meta: {match['meta_file']}")

        if not force:
            if not is_current:
                print("\n⚠️  CROSS-PROJECT REMOVAL!")
                print("   This will delete the index from another project.")

            response = input(f"   ❓ Confirm removal from {location_info}? (y/N): ").strip().lower()
            if response not in ["y", "yes"]:
                print("   ❌ Removal cancelled.")
                return False

        if kind == "cli":
            return self._delete_index_directory(
                match["index_dir"],
                index_name,
                project_path if not is_current else None,
                is_app=False,
            )
        else:
            return self._delete_index_directory(
                match["index_dir"],
                match.get("display_name", index_name),
                project_path if not is_current else None,
                is_app=True,
                meta_file=match.get("meta_file"),
                app_file_base=match.get("file_base"),
            )

    def _remove_from_multiple_matches(self, matches, index_name: str, force: bool):
        """Handle removal when multiple matches are found"""

        print(f"⚠️  Found {len(matches)} indexes named '{index_name}':")
        print("   " + "─" * 50)

        for i, match in enumerate(matches, 1):
            project_path = match["project_path"]
            is_current = match["is_current"]
            kind = match.get("kind", "cli")

            if is_current:
                print(f"   {i}. 🏠 Current project ({'CLI' if kind == 'cli' else 'APP'})")
            else:
                print(f"   {i}. 📂 {project_path.name} ({'CLI' if kind == 'cli' else 'APP'})")

            # Show path details
            if kind == "cli":
                print(f"      📍 {match['path']}")
            else:
                print(f"      📍 {match['meta_file']}")

            # Show size info
            try:
                if kind == "cli":
                    size_mb = sum(
                        f.stat().st_size for f in match["index_dir"].iterdir() if f.is_file()
                    ) / (1024 * 1024)
                else:
                    file_base = match.get("file_base")
                    size_mb = 0.0
                    if file_base:
                        size_mb = sum(
                            f.stat().st_size
                            for f in match["index_dir"].glob(f"{file_base}.leann*")
                            if f.is_file()
                        ) / (1024 * 1024)
                print(f"      📦 Size: {size_mb:.1f} MB")
            except (OSError, PermissionError):
                pass

        print("   " + "─" * 50)

        if force:
            # If all matches point to the same physical location (should be handled by deduplication, but just in case)
            # Or if the user really wants to force delete everything with this name.
            # But for safety, we only allow force delete if there's no ambiguity.
            # Since we added deduplication, multiple matches now means DIFFERENT physical locations.
            print("   ❌ Multiple matches found at DIFFERENT locations, but --force specified.")
            print("   Please run without --force to choose which one to remove safely.")
            return False

        try:
            choice = input(
                f"   ❓ Which one to remove? (1-{len(matches)}, or 'c' to cancel): "
            ).strip()
            if choice.lower() == "c":
                print("   ❌ Removal cancelled.")
                return False

            choice_idx = int(choice) - 1
            if 0 <= choice_idx < len(matches):
                selected_match = matches[choice_idx]
                project_path = selected_match["project_path"]
                is_current = selected_match["is_current"]
                kind = selected_match.get("kind", "cli")

                location = "current project" if is_current else f"'{project_path.name}' project"
                print(f"   🎯 Selected: Remove from {location}")

                # Final confirmation for safety
                confirm = input(
                    f"   ❓ FINAL CONFIRMATION - Type '{index_name}' to proceed: "
                ).strip()
                if confirm != index_name:
                    print("   ❌ Confirmation failed. Removal cancelled.")
                    return False

                if kind == "cli":
                    return self._delete_index_directory(
                        selected_match["index_dir"],
                        index_name,
                        project_path if not is_current else None,
                        is_app=False,
                    )
                else:
                    return self._delete_index_directory(
                        selected_match["index_dir"],
                        selected_match.get("display_name", index_name),
                        project_path if not is_current else None,
                        is_app=True,
                        meta_file=selected_match.get("meta_file"),
                        app_file_base=selected_match.get("file_base"),
                    )
            else:
                print("   ❌ Invalid choice. Removal cancelled.")
                return False

        except (ValueError, KeyboardInterrupt):
            print("\n   ❌ Invalid input. Removal cancelled.")
            return False

    def _delete_index_directory(
        self,
        index_dir: Path,
        index_display_name: str,
        project_path: Optional[Path] = None,
        is_app: bool = False,
        meta_file: Optional[Path] = None,
        app_file_base: Optional[str] = None,
    ):
        """Delete a CLI index directory or APP index files safely."""
        try:
            if is_app:
                removed = 0
                errors = 0
                # Delete only files that belong to this app index (based on file base)
                pattern_base = app_file_base or ""
                for f in index_dir.glob(f"{pattern_base}.leann*"):
                    try:
                        f.unlink()
                        removed += 1
                    except Exception:
                        errors += 1
                # Best-effort: also remove the meta file if specified and still exists
                if meta_file and meta_file.exists():
                    try:
                        meta_file.unlink()
                        removed += 1
                    except Exception:
                        errors += 1

                if removed > 0 and errors == 0:
                    if project_path:
                        print(
                            f"✅ App index '{index_display_name}' removed from {project_path.name}"
                        )
                    else:
                        print(f"✅ App index '{index_display_name}' removed successfully")
                    return True
                elif removed > 0 and errors > 0:
                    print(
                        f"⚠️  App index '{index_display_name}' partially removed (some files couldn't be deleted)"
                    )
                    return True
                else:
                    print(
                        f"❌ No files found to remove for app index '{index_display_name}' in {index_dir}"
                    )
                    return False
            else:
                import shutil

                shutil.rmtree(index_dir)

                if project_path:
                    print(f"✅ Index '{index_display_name}' removed from {project_path.name}")
                else:
                    print(f"✅ Index '{index_display_name}' removed successfully")
                return True
        except Exception as e:
            print(f"❌ Error removing index '{index_display_name}': {e}")
            return False

    def load_documents(
        self,
        docs_paths: Union[str, list],
        custom_file_types: Union[str, None] = None,
        include_hidden: bool = False,
        args: Optional[dict[str, Any]] = None,
    ):
        # Handle both single path (string) and multiple paths (list) for backward compatibility
        if isinstance(docs_paths, str):
            docs_paths = [docs_paths]

        # Separate files and directories
        files = []
        directories = []
        for path in docs_paths:
            path_obj = Path(path)
            if path_obj.is_file():
                files.append(str(path_obj))
            elif path_obj.is_dir():
                # Check if this is a git submodule - if so, skip it
                if self._is_git_submodule(path_obj):
                    print(f"⚠️  Skipping git submodule: {path}")
                    continue
                directories.append(str(path_obj))
            else:
                print(f"⚠️  Warning: Path '{path}' does not exist, skipping...")
                continue

        # Define custom extractors map
        custom_extractors = {
            ".docx": extract_docx_text,
            ".pptx": extract_pptx_text,
            ".xlsx": extract_xlsx_text,
            ".mm": extract_mm_text,
        }

        # Print summary of what we're processing
        total_items = len(files) + len(directories)
        items_desc = []
        if files:
            items_desc.append(f"{len(files)} file{'s' if len(files) > 1 else ''}")
        if directories:
            items_desc.append(
                f"{len(directories)} director{'ies' if len(directories) > 1 else 'y'}"
            )

        print(f"Loading documents from {' and '.join(items_desc)} ({total_items} total):")
        if files:
            print(f"  📄 Files: {', '.join([Path(f).name for f in files])}")
        if directories:
            print(f"  📁 Directories: {', '.join(directories)}")

        if custom_file_types:
            print(f"Using custom file types: {custom_file_types}")

        all_documents = []

        # Helper to detect hidden path components
        def _path_has_hidden_segment(p: Path) -> bool:
            return any(part.startswith(".") and part not in [".", ".."] for part in p.parts)

        # First, process individual files if any
        if files:
            print(f"\n🔄 Processing {len(files)} individual file{'s' if len(files) > 1 else ''}...")

            for file_path in files:
                file_path_obj = Path(file_path)
                if not include_hidden and _path_has_hidden_segment(file_path_obj):
                    print(f"  ⚠️  Skipping hidden file: {file_path}")
                    continue

                ext = file_path_obj.suffix.lower()
                doc = None

                # Try PDF extraction first if it's a PDF
                if ext == ".pdf":
                    text = extract_pdf_text(file_path)
                    if text:
                        from llama_index.core import Document

                        doc = Document(text=text, metadata={"source": file_path})

                # Try custom extractors (DOCX, PPTX, XLSX, MM)
                elif ext in custom_extractors:
                    text = custom_extractors[ext](file_path)
                    if text:
                        from llama_index.core import Document

                        doc = Document(text=text, metadata={"source": file_path})

                # If handled by custom logic, add it
                if doc:
                    all_documents.append(doc)
                    print(f"  ✅ Loaded custom format: {file_path_obj.name}")
                else:
                    # Fallback to SimpleDirectoryReader for other files
                    try:
                        file_docs = SimpleDirectoryReader(
                            str(file_path_obj.parent),
                            input_files=[file_path],
                            filename_as_id=True,
                        ).load_data()
                        for d in file_docs:
                            d.metadata["source"] = file_path
                        all_documents.extend(file_docs)
                        print(f"  ✅ Loaded: {file_path_obj.name}")
                    except Exception as e:
                        print(f"  ❌ Warning: Could not load file {file_path}: {e}")

        # Define file extensions to process
        if custom_file_types:
            # Parse custom file types from comma-separated string
            code_extensions = [ext.strip() for ext in custom_file_types.split(",") if ext.strip()]
            # Ensure extensions start with a dot
            code_extensions = [ext if ext.startswith(".") else f".{ext}" for ext in code_extensions]
        else:
            # Use default supported file types
            code_extensions = [
                ".txt",
                ".md",
                ".docx",
                ".pptx",
                ".xlsx",
                ".mm",
                ".py",
                ".js",
                ".ts",
                ".jsx",
                ".tsx",
                ".java",
                ".cpp",
                ".c",
                ".h",
                ".hpp",
                ".cs",
                ".go",
                ".rs",
                ".rb",
                ".php",
                ".swift",
                ".kt",
                ".scala",
                ".r",
                ".sql",
                ".sh",
                ".bash",
                ".zsh",
                ".fish",
                ".ps1",
                ".bat",
                # Config and markup files
                ".json",
                ".yaml",
                ".yml",
                ".xml",
                ".toml",
                ".ini",
                ".cfg",
                ".conf",
                ".html",
                ".css",
                ".scss",
                ".less",
                ".vue",
                ".svelte",
                # Data science
                ".ipynb",
                ".R",
                ".py",
                ".jl",
            ]

        # Process each directory
        if directories:
            print(
                f"\n🔄 Processing {len(directories)} director{'ies' if len(directories) > 1 else 'y'}..."
            )

        for docs_dir in directories:
            print(f"Processing directory: {docs_dir}")
            # Build gitignore parser for each directory
            gitignore_matches = self._build_gitignore_parser(docs_dir)

            # Try to use better PDF parsers first, but only if PDFs are requested
            documents = []
            # Use resolved absolute paths to avoid mismatches (symlinks, relative vs absolute)
            docs_path = Path(docs_dir).resolve()

            # Check if we should process PDFs
            should_process_pdfs = custom_file_types is None or ".pdf" in custom_file_types

            if should_process_pdfs:
                for file_path in docs_path.rglob("*.pdf"):
                    # Check if file matches any exclude pattern
                    try:
                        # Ensure both paths are resolved before computing relativity
                        file_path_resolved = file_path.resolve()
                        # Determine directory scope using the non-resolved path to avoid
                        # misclassifying symlinked entries as outside the docs directory
                        relative_path = file_path.relative_to(docs_path)
                        if not include_hidden and _path_has_hidden_segment(relative_path):
                            continue
                        # Use absolute path for gitignore matching
                        if self._should_exclude_file(file_path_resolved, gitignore_matches):
                            continue
                    except ValueError:
                        # Skip files that can't be made relative to docs_path
                        print(f"⚠️  Skipping file outside directory scope: {file_path}")
                        continue

                    print(f"Processing PDF: {file_path}")

                    text = extract_pdf_text(str(file_path))

                    if text:
                        # Create a simple document structure
                        from llama_index.core import Document

                        doc = Document(text=text, metadata={"source": str(file_path)})
                        documents.append(doc)
                    else:
                        # Fallback to default reader
                        print(f"Using default reader for {file_path}")
                        try:
                            default_docs = SimpleDirectoryReader(
                                str(file_path.parent),
                                exclude_hidden=not include_hidden,
                                filename_as_id=True,
                                required_exts=[file_path.suffix],
                            ).load_data()
                            documents.extend(default_docs)
                        except Exception as e:
                            print(f"Warning: Could not process {file_path}: {e}")

            for ext, extractor in custom_extractors.items():
                should_process = custom_file_types is None or ext in custom_file_types
                if not should_process:
                    continue
                for file_path in docs_path.rglob(f"*{ext}"):
                    try:
                        file_path_resolved = file_path.resolve()
                        relative_path = file_path.relative_to(docs_path)
                        if not include_hidden and _path_has_hidden_segment(relative_path):
                            continue
                        if self._should_exclude_file(file_path_resolved, gitignore_matches):
                            continue
                    except ValueError:
                        print(f"⚠️  Skipping file outside directory scope: {file_path}")
                        continue
                    print(f"Processing {ext.upper()[1:]}: {file_path}")
                    text = extractor(str(file_path))
                    if text:
                        from llama_index.core import Document

                        doc = Document(text=text, metadata={"source": str(file_path)})
                        documents.append(doc)
                    else:
                        print(f"⚠️  Could not extract text from: {file_path}")

            other_file_extensions = code_extensions
            processed_extensions = [".pdf", *custom_extractors.keys()]
            other_file_extensions = [
                ext for ext in code_extensions if ext not in processed_extensions
            ]

            try:

                def is_valid_file(file_path: Path) -> bool:
                    try:
                        if file_path.is_symlink():
                            resolved = file_path.resolve()
                            if not resolved.exists():
                                return False
                        elif not file_path.exists():
                            return False
                        return True
                    except (OSError, PermissionError):
                        return False

                def file_filter(
                    file_path: str, docs_dir=docs_dir, gitignore_matches=gitignore_matches
                ) -> bool:
                    try:
                        file_path_obj = Path(file_path)
                        if not is_valid_file(file_path_obj):
                            return False
                        docs_path_obj = Path(docs_dir).resolve()
                        file_path_obj = file_path_obj.resolve()
                        _ = file_path_obj.relative_to(docs_path_obj)
                        return not self._should_exclude_file(file_path_obj, gitignore_matches)
                    except (ValueError, OSError):
                        return False

                if other_file_extensions:
                    try:
                        docs_path = Path(docs_dir)
                        valid_files = []
                        ext_set = set(other_file_extensions)
                        for file_path in docs_path.rglob("*"):
                            if file_path.suffix.lower() not in ext_set:
                                continue
                            if not include_hidden and _path_has_hidden_segment(
                                file_path.relative_to(docs_path)
                            ):
                                continue
                            if not is_valid_file(file_path):
                                print(f"⚠️  Skipping broken symlink: {file_path}")
                                continue
                            if self._should_exclude_file(file_path.resolve(), gitignore_matches):
                                continue
                            valid_files.append(str(file_path))

                        if valid_files:
                            other_docs = SimpleDirectoryReader(
                                input_files=valid_files,
                                encoding="utf-8",
                                file_extractor={},
                                filename_as_id=True,
                            ).load_data(show_progress=True)
                        else:
                            other_docs = []
                    except FileNotFoundError as fnf_err:
                        print(f"⚠️  Skipping broken symlink or missing file: {fnf_err}")
                        other_docs = []
                    except Exception as load_err:
                        print(f"⚠️  Error loading documents from {docs_dir}: {load_err}")
                        other_docs = []
                else:
                    other_docs = []

                for doc in other_docs:
                    try:
                        file_path = doc.metadata.get("file_path", "")
                        doc.metadata["source"] = file_path
                        documents.append(doc)
                    except Exception as doc_err:
                        print(f"⚠️  Skipping document due to metadata error: {doc_err}")
                        continue

            except ValueError as e:
                if "No files found" in str(e):
                    print(f"No additional files found for other supported types in {docs_dir}.")
                else:
                    raise e

            all_documents.extend(documents)
            print(f"Loaded {len(documents)} documents from {docs_dir}")

        documents = all_documents

        all_texts = []

        # Define code file extensions for intelligent chunking
        code_file_exts = {
            ".py",
            ".js",
            ".ts",
            ".jsx",
            ".tsx",
            ".java",
            ".cpp",
            ".c",
            ".h",
            ".hpp",
            ".cs",
            ".go",
            ".rs",
            ".rb",
            ".php",
            ".swift",
            ".kt",
            ".scala",
            ".r",
            ".sql",
            ".sh",
            ".bash",
            ".zsh",
            ".fish",
            ".ps1",
            ".bat",
            ".json",
            ".yaml",
            ".yml",
            ".xml",
            ".toml",
            ".ini",
            ".cfg",
            ".conf",
            ".html",
            ".css",
            ".scss",
            ".less",
            ".vue",
            ".svelte",
            ".ipynb",
            ".R",
            ".jl",
        }

        print("start chunking documents")

        # Check if AST chunking is requested
        use_ast = getattr(args, "use_ast_chunking", False)

        if use_ast:
            print("🧠 Using AST-aware chunking for code files")
            try:
                # Import enhanced chunking utilities from packaged module
                from .chunking_utils import create_text_chunks

                # Use enhanced chunking with AST support
                chunk_texts = create_text_chunks(
                    documents,
                    chunk_size=self.node_parser.chunk_size,
                    chunk_overlap=self.node_parser.chunk_overlap,
                    use_ast_chunking=True,
                    ast_chunk_size=getattr(args, "ast_chunk_size", 768),
                    ast_chunk_overlap=getattr(args, "ast_chunk_overlap", 96),
                    code_file_extensions=None,  # Use defaults
                    ast_fallback_traditional=getattr(args, "ast_fallback_traditional", True),
                )

                # create_text_chunks now returns list[dict] with metadata preserved
                all_texts.extend(chunk_texts)

            except ImportError as e:
                print(
                    f"⚠️  AST chunking utilities not available in package ({e}), falling back to traditional chunking"
                )
                use_ast = False

        if not use_ast:
            # Use traditional chunking logic
            for doc in tqdm(documents, desc="Chunking documents", unit="doc"):
                # Check if this is a code file based on source path
                source_path = doc.metadata.get("source", "")
                file_path = doc.metadata.get("file_path", "")
                is_code_file = any(source_path.endswith(ext) for ext in code_file_exts)

                # Extract metadata to preserve with chunks
                chunk_metadata = {
                    "file_path": file_path or source_path,
                    "file_name": doc.metadata.get("file_name", ""),
                    "source": source_path,
                }

                # Add optional metadata if available
                if "creation_date" in doc.metadata:
                    chunk_metadata["creation_date"] = doc.metadata["creation_date"]
                if "last_modified_date" in doc.metadata:
                    chunk_metadata["last_modified_date"] = doc.metadata["last_modified_date"]

                # Use appropriate parser based on file type
                parser = self.code_parser if is_code_file else self.node_parser
                nodes = parser.get_nodes_from_documents([doc])

                for node in nodes:
                    all_texts.append({"text": node.get_content(), "metadata": chunk_metadata})

        print(f"Loaded {len(documents)} documents, {len(all_texts)} chunks")
        return all_texts

    async def build_index(self, args):
        docs_paths = args.docs
        # Use current directory name if index_name not provided
        if args.index_name:
            index_name = args.index_name
        else:
            index_name = Path.cwd().name
            print(f"Using current directory name as index: '{index_name}'")

        index_dir = self.indexes_dir / index_name
        index_path = self.get_index_path(index_name)

        # Display all paths being indexed with file/directory distinction
        files = [p for p in docs_paths if Path(p).is_file()]
        directories = [p for p in docs_paths if Path(p).is_dir()]

        print(f"📂 Indexing {len(docs_paths)} path{'s' if len(docs_paths) > 1 else ''}:")
        if files:
            print(f"  📄 Files ({len(files)}):")
            for i, file_path in enumerate(files, 1):
                print(f"    {i}. {Path(file_path).resolve()}")
        if directories:
            print(f"  📁 Directories ({len(directories)}):")
            for i, dir_path in enumerate(directories, 1):
                print(f"    {i}. {Path(dir_path).resolve()}")

        if index_dir.exists() and not args.force:
            print(f"Index '{index_name}' already exists. Use --force to rebuild.")
            return

        # Configure chunking based on CLI args before loading documents
        # Guard against invalid configurations
        doc_chunk_size = max(1, int(args.doc_chunk_size))
        doc_chunk_overlap = max(0, int(args.doc_chunk_overlap))
        if doc_chunk_overlap >= doc_chunk_size:
            print(
                f"⚠️  Adjusting doc chunk overlap from {doc_chunk_overlap} to {doc_chunk_size - 1} (must be < chunk size)"
            )
            doc_chunk_overlap = doc_chunk_size - 1

        code_chunk_size = max(1, int(args.code_chunk_size))
        code_chunk_overlap = max(0, int(args.code_chunk_overlap))
        if code_chunk_overlap >= code_chunk_size:
            print(
                f"⚠️  Adjusting code chunk overlap from {code_chunk_overlap} to {code_chunk_size - 1} (must be < chunk size)"
            )
            code_chunk_overlap = code_chunk_size - 1

        self.node_parser = SentenceSplitter(
            chunk_size=doc_chunk_size,
            chunk_overlap=doc_chunk_overlap,
            separator=" ",
            paragraph_separator="\n\n",
        )
        self.code_parser = SentenceSplitter(
            chunk_size=code_chunk_size,
            chunk_overlap=code_chunk_overlap,
            separator="\n",
            paragraph_separator="\n\n",
        )

        all_texts = self.load_documents(
            docs_paths, args.file_types, include_hidden=args.include_hidden, args=args
        )
        if not all_texts:
            print("No documents found")
            return

        index_dir.mkdir(parents=True, exist_ok=True)

        print(f"Building index '{index_name}' with {args.backend_name} backend...")

        embedding_options: dict[str, Any] = {}
        if args.embedding_mode == "ollama":
            embedding_options["host"] = resolve_ollama_host(args.embedding_host)
        elif args.embedding_mode == "openai":
            embedding_options["base_url"] = resolve_openai_base_url(args.embedding_api_base)
            resolved_embedding_key = resolve_openai_api_key(args.embedding_api_key)
            if resolved_embedding_key:
                embedding_options["api_key"] = resolved_embedding_key
        if args.query_prompt_template:
            # New format: separate templates
            if args.embedding_prompt_template:
                embedding_options["build_prompt_template"] = args.embedding_prompt_template
            embedding_options["query_prompt_template"] = args.query_prompt_template
        elif args.embedding_prompt_template:
            # Old format: single template (backward compat)
            embedding_options["prompt_template"] = args.embedding_prompt_template

        builder = LeannBuilder(
            backend_name=args.backend_name,
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            embedding_options=embedding_options or None,
            graph_degree=args.graph_degree,
            complexity=args.complexity,
            is_compact=args.compact,
            is_recompute=args.recompute,
            num_threads=args.num_threads,
        )

        for chunk in all_texts:
            builder.add_text(chunk["text"], metadata=chunk["metadata"])

        builder.build_index(index_path)
        print(f"Index built at {index_path}")

        # Register this project directory in global registry
        self.register_project_dir()

    async def index_browser(self, args):
        """Build an index from browser history."""
        from .readers import ChromeHistoryReader

        browser_type = args.browser_type
        profile = args.profile
        index_name = args.index_name or f"{browser_type}-history"

        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print(f"🌐 Indexing {browser_type.capitalize()} history (profile: {profile})...")

        # Find browser path
        paths = ChromeHistoryReader.find_browser_paths()
        if browser_type not in paths:
            print(f"❌ Could not find {browser_type} profile directory automatically.")
            return

        profile_path = paths[browser_type] / profile

        reader = ChromeHistoryReader()
        documents = reader.load_data(
            chrome_profile_path=str(profile_path), max_count=args.max_items
        )

        if not documents:
            print("❌ No history entries found to index.")
            return

        print(f"📚 Loaded {len(documents)} entries. Building index...")

        index_dir.mkdir(parents=True, exist_ok=True)

        embedding_options = {}
        if args.embedding_mode == "ollama":
            embedding_options["host"] = resolve_ollama_host(None)

        builder = LeannBuilder(
            backend_name="hnsw",  # Default to hnsw for browser history
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            embedding_options=embedding_options or None,
            is_recompute=False,  # Store embeddings for fast browser history search
            is_compact=False,  # Disable compact storage to keep full embeddings
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ Browser history index built at: {index_path}")
        print(f'   Usage: leann search {index_name} "query"')

    async def index_email(self, args):
        """Build an index from Apple Mail emails."""
        from .readers import AppleMailReader

        index_name = args.index_name
        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print("📧 Indexing Apple Mail emails...")

        reader = AppleMailReader()
        documents = reader.load_data(max_count=args.max_items)

        if not documents:
            print("❌ No emails found to index. Make sure Full Disk Access is granted.")
            return

        print(f"📚 Loaded {len(documents)} emails. Building index...")
        index_dir.mkdir(parents=True, exist_ok=True)

        embedding_options = {}
        if args.embedding_mode == "ollama":
            embedding_options["host"] = resolve_ollama_host(None)

        builder = LeannBuilder(
            backend_name="hnsw",
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            embedding_options=embedding_options or None,
            is_recompute=False,
            is_compact=False,
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ Email index built at: {index_path}")
        print(f'   Usage: leann search {index_name} "query"')

    async def index_calendar(self, args):
        """Build an index from Apple Calendar events."""
        from .readers import AppleCalendarReader

        index_name = args.index_name
        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print("📅 Indexing Apple Calendar events...")

        reader = AppleCalendarReader()
        documents = reader.load_data(max_count=args.max_items)

        if not documents:
            print("❌ No calendar events found to index. Make sure Full Disk Access is granted.")
            return

        print(f"📚 Loaded {len(documents)} events. Building index...")
        index_dir.mkdir(parents=True, exist_ok=True)

        embedding_options = {}
        if args.embedding_mode == "ollama":
            embedding_options["host"] = resolve_ollama_host(None)

        builder = LeannBuilder(
            backend_name="hnsw",
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            embedding_options=embedding_options or None,
            is_recompute=False,
            is_compact=False,
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ Calendar index built at: {index_path}")
        print(f'   Usage: leann search {index_name} "query"')

    async def index_wechat(self, args):
        """Build an index from WeChat chat history."""
        from .readers import WeChatHistoryReader

        index_name = args.index_name
        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print(f"💬 Indexing WeChat history from {args.export_dir}...")

        reader = WeChatHistoryReader()
        documents = reader.load_data(wechat_export_dir=args.export_dir, max_count=args.max_items)

        if not documents:
            print("❌ No WeChat history found.")
            return

        print(f"📚 Loaded {len(documents)} chat threads. Building index...")
        index_dir.mkdir(parents=True, exist_ok=True)

        builder = LeannBuilder(
            backend_name="hnsw",
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            is_recompute=False,
            is_compact=False,
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ WeChat index built at: {index_path}")

    async def index_imessage(self, args):
        """Build an index from iMessage history."""
        from .readers import IMessageReader

        index_name = args.index_name
        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print("💬 Indexing iMessage history...")

        reader = IMessageReader(concatenate_conversations=True)
        documents = reader.load_data(input_dir=args.db_path)

        if not documents:
            print("❌ No iMessage history found.")
            return

        if args.max_items > 0:
            documents = documents[: args.max_items]

        print(f"📚 Loaded {len(documents)} conversations. Building index...")
        index_dir.mkdir(parents=True, exist_ok=True)

        builder = LeannBuilder(
            backend_name="hnsw",
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            is_recompute=False,
            is_compact=False,
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ iMessage index built at: {index_path}")

    async def index_slack(self, args):
        """Build an index from Slack via MCP."""
        from .readers import SlackMCPReader

        index_name = args.index_name
        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print(f"💬 Indexing Slack via {args.mcp_server}...")

        reader = SlackMCPReader(
            mcp_server_command=args.mcp_server,
            workspace_name=args.workspace,
            concatenate_conversations=True,
        )
        documents = await reader.load_data(channels=args.channels)

        if not documents:
            print("❌ No Slack messages found or MCP server failed.")
            return

        print(f"📚 Loaded {len(documents)} items. Building index...")
        index_dir.mkdir(parents=True, exist_ok=True)

        builder = LeannBuilder(
            backend_name="hnsw",
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            is_recompute=False,
            is_compact=False,
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ Slack index built at: {index_path}")

    async def index_chatgpt(self, args):
        """Build an index from ChatGPT export."""
        from .readers import ChatGPTReader

        index_name = args.index_name
        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print(f"🤖 Indexing ChatGPT export from {args.export_path}...")

        reader = ChatGPTReader()
        documents = reader.load_data(export_path=args.export_path)

        if not documents:
            print("❌ No ChatGPT data found.")
            return

        print(f"📚 Loaded {len(documents)} conversations. Building index...")
        index_dir.mkdir(parents=True, exist_ok=True)

        builder = LeannBuilder(
            backend_name="hnsw",
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            is_recompute=False,
            is_compact=False,
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ ChatGPT index built at: {index_path}")

    async def index_claude(self, args):
        """Build an index from Claude export."""
        from .readers import ClaudeReader

        index_name = args.index_name
        index_dir = self.indexes_dir / index_name
        index_path = str(index_dir / "documents.leann")

        print(f"🤖 Indexing Claude export from {args.export_path}...")

        reader = ClaudeReader()
        documents = reader.load_data(export_path=args.export_path)

        if not documents:
            print("❌ No Claude data found.")
            return

        print(f"📚 Loaded {len(documents)} conversations. Building index...")
        index_dir.mkdir(parents=True, exist_ok=True)

        builder = LeannBuilder(
            backend_name="hnsw",
            embedding_model=args.embedding_model,
            embedding_mode=args.embedding_mode,
            is_recompute=False,
            is_compact=False,
        )

        for doc in documents:
            builder.add_text(doc.text, metadata=doc.metadata)

        builder.build_index(index_path)
        print(f"✅ Claude index built at: {index_path}")

    async def search_documents(self, args):
        index_name = args.index_name
        query = args.query

        # First try to find the index in current project
        index_path = self.get_index_path(index_name)
        if self.index_exists(index_name):
            # Found in current project, use it
            pass
        else:
            # Search across all registered projects (like list_indexes does)
            all_matches = self._find_all_matching_indexes(index_name)
            if not all_matches:
                print(
                    f"Index '{index_name}' not found. Use 'leann build {index_name} --docs <dir> [<dir2> ...]' to create it."
                )
                return
            elif len(all_matches) == 1:
                # Found exactly one match, use it
                match = all_matches[0]
                if match["kind"] == "cli":
                    index_path = str(match["index_dir"] / "documents.leann")
                else:
                    # App format: use the meta file to construct the path
                    meta_file = match["meta_file"]
                    file_base = match["file_base"]
                    index_path = str(meta_file.parent / f"{file_base}.leann")

                project_info = (
                    "current project"
                    if match["is_current"]
                    else f"project '{match['project_path'].name}'"
                )
                print(f"Using index '{index_name}' from {project_info}")
            else:
                # Multiple matches found
                if args.non_interactive:
                    # Non-interactive mode: automatically select the best match
                    # Priority: current project first, then first available
                    current_matches = [m for m in all_matches if m["is_current"]]
                    if current_matches:
                        match = current_matches[0]
                        location_desc = "current project"
                    else:
                        match = all_matches[0]
                        location_desc = f"project '{match['project_path'].name}'"

                    if match["kind"] == "cli":
                        index_path = str(match["index_dir"] / "documents.leann")
                    else:
                        meta_file = match["meta_file"]
                        file_base = match["file_base"]
                        index_path = str(meta_file.parent / f"{file_base}.leann")

                    print(
                        f"Found {len(all_matches)} indexes named '{index_name}', using index from {location_desc}"
                    )
                else:
                    # Interactive mode: ask user to choose
                    print(f"Found {len(all_matches)} indexes named '{index_name}':")
                    for i, match in enumerate(all_matches, 1):
                        project_path = match["project_path"]
                        is_current = match["is_current"]
                        kind = match.get("kind", "cli")

                        if is_current:
                            print(
                                f"   {i}. 🏠 Current project ({'CLI' if kind == 'cli' else 'APP'})"
                            )
                        else:
                            print(
                                f"   {i}. 📂 {project_path.name} ({'CLI' if kind == 'cli' else 'APP'})"
                            )

                    try:
                        choice = input(f"Which index to search? (1-{len(all_matches)}): ").strip()
                        choice_idx = int(choice) - 1
                        if 0 <= choice_idx < len(all_matches):
                            match = all_matches[choice_idx]
                            if match["kind"] == "cli":
                                index_path = str(match["index_dir"] / "documents.leann")
                            else:
                                meta_file = match["meta_file"]
                                file_base = match["file_base"]
                                index_path = str(meta_file.parent / f"{file_base}.leann")

                            project_info = (
                                "current project"
                                if match["is_current"]
                                else f"project '{match['project_path'].name}'"
                            )
                            print(f"Using index '{index_name}' from {project_info}")
                        else:
                            print("Invalid choice. Aborting search.")
                            return
                    except (ValueError, KeyboardInterrupt):
                        print("Invalid input. Aborting search.")
                        return

        # Build provider_options for runtime override
        provider_options = {}
        if args.embedding_prompt_template:
            provider_options["prompt_template"] = args.embedding_prompt_template

        searcher = LeannSearcher(index_path=index_path)
        results = searcher.search(
            query,
            top_k=args.top_k,
            complexity=args.complexity,
            beam_width=args.beam_width,
            prune_ratio=args.prune_ratio,
            recompute_embeddings=args.recompute_embeddings,
            pruning_strategy=args.pruning_strategy,
            provider_options=provider_options if provider_options else None,
        )

        print(f"Search results for '{query}' (top {len(results)}):")
        for i, result in enumerate(results, 1):
            print(f"{i}. Score: {result.score:.3f}")

            # Display metadata if flag is set
            if args.show_metadata and result.metadata:
                file_path = result.metadata.get("file_path", "")
                if file_path:
                    print(f"   📄 File: {file_path}")

                file_name = result.metadata.get("file_name", "")
                if file_name and file_name != file_path:
                    print(f"   📝 Name: {file_name}")

                # Show timestamps if available
                if "creation_date" in result.metadata:
                    print(f"   🕐 Created: {result.metadata['creation_date']}")
                if "last_modified_date" in result.metadata:
                    print(f"   🕑 Modified: {result.metadata['last_modified_date']}")

            print(f"   {result.text[:200]}...")
            print(f"   Source: {result.metadata.get('source', '')}")
            print()

    async def ask_questions(self, args):
        index_name = args.index_name
        index_path = self.get_index_path(index_name)

        if not self.index_exists(index_name):
            print(
                f"Index '{index_name}' not found. Use 'leann build {index_name} --docs <dir> [<dir2> ...]' to create it."
            )
            return

        print(f"Starting chat with index '{index_name}'...")
        print(f"Using {args.model} ({args.llm})")

        llm_config = {"type": args.llm, "model": args.model}
        if args.llm == "ollama":
            llm_config["host"] = resolve_ollama_host(args.host)
        elif args.llm == "openai":
            llm_config["base_url"] = resolve_openai_base_url(args.api_base)
            resolved_api_key = resolve_openai_api_key(args.api_key)
            if resolved_api_key:
                llm_config["api_key"] = resolved_api_key
        elif args.llm == "anthropic":
            # For Anthropic, pass base_url and API key if provided
            if args.api_base:
                llm_config["base_url"] = resolve_anthropic_base_url(args.api_base)
            if args.api_key:
                llm_config["api_key"] = args.api_key

        chat = LeannChat(index_path=index_path, llm_config=llm_config)

        llm_kwargs: dict[str, Any] = {}
        if args.thinking_budget:
            llm_kwargs["thinking_budget"] = args.thinking_budget

        def _ask_once(prompt: str) -> None:
            query_start_time = time.time()
            response = chat.ask(
                prompt,
                top_k=args.top_k,
                complexity=args.complexity,
                beam_width=args.beam_width,
                prune_ratio=args.prune_ratio,
                recompute_embeddings=args.recompute_embeddings,
                pruning_strategy=args.pruning_strategy,
                llm_kwargs=llm_kwargs,
            )
            query_completion_time = time.time() - query_start_time
            print(f"LEANN: {response}")
            print(f"The query took {query_completion_time:.3f} seconds to finish")

        initial_query = (args.query or "").strip()

        if args.interactive:
            # Create interactive session
            session = create_cli_session(index_name)

            if initial_query:
                _ask_once(initial_query)

            session.run_interactive_loop(_ask_once)
        else:
            query = initial_query or input("Enter your question: ").strip()
            if not query:
                print("No question provided. Exiting.")
                return

            _ask_once(query)

    async def react_agent(self, args):
        """Run ReAct agent for multiturn retrieval."""
        index_name = args.index_name
        query = args.query

        # Find the index (similar to search_documents)
        index_path = self.get_index_path(index_name)
        if self.index_exists(index_name):
            pass
        else:
            all_matches = self._find_all_matching_indexes(index_name)
            if not all_matches:
                print(
                    f"Index '{index_name}' not found. Use 'leann build {index_name} --docs <dir> [<dir2> ...]' to create it."
                )
                return
            elif len(all_matches) == 1:
                match = all_matches[0]
                if match["kind"] == "cli":
                    index_path = str(match["index_dir"] / "documents.leann")
                else:
                    meta_file = match["meta_file"]
                    file_base = match["file_base"]
                    index_path = str(meta_file.parent / f"{file_base}.leann")
            else:
                # Multiple matches - use first one for now
                match = all_matches[0]
                if match["kind"] == "cli":
                    index_path = str(match["index_dir"] / "documents.leann")
                else:
                    meta_file = match["meta_file"]
                    file_base = match["file_base"]
                    index_path = str(meta_file.parent / f"{file_base}.leann")
                print(f"Found {len(all_matches)} indexes named '{index_name}', using first match")

        print(f"🤖 Starting ReAct agent with index '{index_name}'...")
        print(f"Using {args.model} ({args.llm})")

        llm_config = {"type": args.llm, "model": args.model}
        if args.llm == "ollama":
            llm_config["host"] = resolve_ollama_host(args.host)
        elif args.llm == "openai":
            llm_config["base_url"] = resolve_openai_base_url(args.api_base)
            resolved_api_key = resolve_openai_api_key(args.api_key)
            if resolved_api_key:
                llm_config["api_key"] = resolved_api_key
        elif args.llm == "anthropic":
            if args.api_base:
                llm_config["base_url"] = resolve_anthropic_base_url(args.api_base)
            if args.api_key:
                llm_config["api_key"] = args.api_key

        from .react_agent import create_react_agent

        agent = create_react_agent(
            index_path=index_path,
            llm_config=llm_config,
            max_iterations=args.max_iterations,
        )

        print(f"\n🔍 Question: {query}\n")
        answer = agent.run(query, top_k=args.top_k)
        print(f"\n✅ Final Answer:\n{answer}\n")

        if agent.search_history:
            print(f"\n📊 Search History ({len(agent.search_history)} iterations):")
            for entry in agent.search_history:
                print(
                    f"  {entry['iteration']}. {entry['action']} ({entry['results_count']} results)"
                )

    async def serve_api(self, args):
        """Start the HTTP API server."""
        import os

        try:
            from .server import main as server_main

            # Override host/port if provided via CLI args
            if args.host:
                os.environ["LEANN_SERVER_HOST"] = args.host
            if args.port:
                os.environ["LEANN_SERVER_PORT"] = str(args.port)

            # Run the server (this is blocking, so we don't await it)
            # The server_main function handles uvicorn.run which blocks
            server_main()
        except ImportError as e:
            print(
                "❌ HTTP server dependencies not installed.\n"
                "Install them with:\n"
                "  uv pip install 'leann-core[server]'\n"
                "or:\n"
                "  uv pip install 'fastapi>=0.115' 'pydantic>=2' 'uvicorn[standard]'\n"
            )
            raise SystemExit(1) from e
        except Exception as e:
            print(f"❌ Error starting server: {e}")
            raise SystemExit(1) from e

    async def run(self, args=None):
        parser = self.create_parser()

        if args is None:
            args = parser.parse_args()

        if not args.command:
            parser.print_help()
            return

        # Determine whether to suppress C++ output
        # Default is to suppress (quiet mode), unless --verbose is specified
        suppress = not getattr(args, "verbose", False)

        if args.command == "list":
            self.list_indexes()
        elif args.command == "remove":
            self.remove_index(args.index_name, args.force)
        elif args.command == "build":
            with suppress_cpp_output(suppress):
                await self.build_index(args)
        elif args.command == "search":
            with suppress_cpp_output(suppress):
                await self.search_documents(args)
        elif args.command == "ask":
            with suppress_cpp_output(suppress):
                await self.ask_questions(args)
        elif args.command == "react":
            with suppress_cpp_output(suppress):
                await self.react_agent(args)
        elif args.command == "index-browser":
            with suppress_cpp_output(suppress):
                await self.index_browser(args)
        elif args.command == "index-email":
            with suppress_cpp_output(suppress):
                await self.index_email(args)
        elif args.command == "index-calendar":
            with suppress_cpp_output(suppress):
                await self.index_calendar(args)
        elif args.command == "index-wechat":
            with suppress_cpp_output(suppress):
                await self.index_wechat(args)
        elif args.command == "index-imessage":
            with suppress_cpp_output(suppress):
                await self.index_imessage(args)
        elif args.command == "index-slack":
            with suppress_cpp_output(suppress):
                await self.index_slack(args)
        elif args.command == "index-chatgpt":
            with suppress_cpp_output(suppress):
                await self.index_chatgpt(args)
        elif args.command == "index-claude":
            with suppress_cpp_output(suppress):
                await self.index_claude(args)
        elif args.command == "serve":
            await self.serve_api(args)
        else:
            parser.print_help()


def main():
    import logging

    import dotenv

    dotenv.load_dotenv()

    # Set clean logging for CLI usage
    logging.getLogger().setLevel(logging.WARNING)  # Only show warnings and errors

    cli = LeannCLI()
    asyncio.run(cli.run())


if __name__ == "__main__":
    main()
